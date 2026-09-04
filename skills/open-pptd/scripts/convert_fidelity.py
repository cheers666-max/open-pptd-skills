#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""convert_fidelity.py — PPTX → PPTD 转换 + 保真度评分。

把现有 PPTX 转成可编辑的 PPTD 项目，并对每个元素给出置信度评分：
  - ≥90%: 自动转换，无需人工复核
  - 70-90%: 转换但标记为"建议复核"
  - <70%: 标记为"需手动重建"，并在报告中给出原因

用法：
  python3 convert_fidelity.py <input.pptx> <output_dir> [--json] [--report report.json]
"""
from __future__ import annotations

import argparse
import json
import math
import os
import sys
import uuid
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches, Cm
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 引擎画布（与 pptd.mjs 保持一致）
ENGINE_W = 1280
ENGINE_H = 720

# OOXML → PPTD 形状映射（覆盖 pptd.mjs 支持的预设）
OOXML_TO_PPTD_SHAPE = {
    "rect": "rect",
    "roundRect": "roundRect",
    "ellipse": "ellipse",
    "triangle": "triangle",
    "diamond": "diamond",
    "parallelogram": "parallelogram",
    "trapezoid": "trapezoid",
    "hexagon": "hexagon",
    "pentagon": "pentagon",
    "octagon": "octagon",
    "rightArrow": "rightArrow",
    "leftArrow": "leftArrow",
    "upArrow": "upArrow",
    "downArrow": "downArrow",
    "chevron": "chevron",
    "homePlate": "homePlate",
    "star5": "star5",
    "star4": "star4",
    "star6": "star6",
    "cross": "cross",
    "donut": "donut",
    "cloud": "cloud",
    "heart": "heart",
    "lightningBolt": "lightningBolt",
    "smileyFace": "smileyFace",
    "sun": "sun",
    "moon": "moon",
    "banner": "banner",
    "ribbon": "ribbon",
    "seal": "seal",
    "gear": "gear",
    "foldedCorner": "foldedCorner",
    "pie": "pie",
    "arc": "arc",
    "blockArc": "blockArc",
}

# 高置信度映射（直接 1:1 转换）
HIGH_CONFIDENCE_SHAPES = {
    "rect", "roundRect", "ellipse", "triangle", "diamond",
    "parallelogram", "trapezoid", "hexagon", "pentagon", "octagon",
    "rightArrow", "leftArrow", "upArrow", "downArrow",
    "chevron", "homePlate", "star5", "star4", "star6",
    "cross", "donut", "cloud", "heart", "lightningBolt",
    "smileyFace", "sun", "moon", "banner", "ribbon", "seal", "gear",
    "foldedCorner", "pie", "arc", "blockArc",
}


def emu_to_px(emu: int) -> float:
    """EMU → px（96 DPI, 1 inch = 914400 EMU）。"""
    return emu / 914400 * 96


def px_to_engine(px_x: float, px_y: float, slide_w_px: float, slide_h_px: float) -> List[float]:
    """Slide 像素坐标 → 引擎坐标（1280x720）。"""
    x = px_x / slide_w_px * ENGINE_W
    y = px_y / slide_h_px * ENGINE_H
    return [round(x, 2), round(y, 2)]


def shape_confidence(shape_type: str, has_text: bool, has_fill: bool, is_group: bool) -> Tuple[float, str]:
    """计算单个元素的置信度并给出原因。"""
    if is_group:
        return 0.50, "组合形状需解组后逐元素处理"

    if shape_type == "picture":
        return 0.95, "图片直接嵌入"

    if shape_type == "text":
        if has_text:
            return 0.92, "文本框直接转换"
        return 0.70, "空文本框（可能含特殊格式）"

    if shape_type == "table":
        return 0.88, "表格结构转换（合并单元格可能丢失）"

    if shape_type == "chart":
        return 0.75, "图表数据可提取，样式需重建"

    if shape_type == "connector" or shape_type == "line":
        return 0.90, "线条/连接线直接转换"

    if shape_type in HIGH_CONFIDENCE_SHAPES:
        return 0.95, f"预设形状 {shape_type} 直接映射"

    if shape_type == "freeform":
        return 0.40, "自由形状需手动重建为 SVG path"

    if shape_type == "ole":
        return 0.30, "OLE 嵌入对象不支持转换"

    if shape_type == "media":
        return 0.35, "音视频媒体需单独处理"

    if shape_type == "smartArt":
        return 0.45, "SmartArt 需解组后重建"

    if shape_type == "diagram":
        return 0.40, "组织结构图/图解需重建"

    return 0.50, f"未知类型 {shape_type}，需人工判断"


def normalize_color(rgb_obj) -> str:
    """python-pptx RGBColor → '#RRGGBB' 格式。"""
    s = str(rgb_obj)
    if not s.startswith("#"):
        s = f"#{s}"
    return s.upper()


def convert_pptx_to_pptd(input_path: Path, output_dir: Path) -> Dict[str, Any]:
    """核心转换：PPTX → PPTD。"""
    if not HAS_PPTX:
        raise RuntimeError("python-pptx 未安装，请运行: pip install python-pptx")

    prs = Presentation(str(input_path))
    # prs.slide_width/height 是 EMU（Length 对象），直接用 EMU 做比例换算
    slide_w = prs.slide_width   # EMU
    slide_h = prs.slide_height  # EMU
    slide_w_px = emu_to_px(slide_w)
    slide_h_px = emu_to_px(slide_h)

    # 创建输出目录
    output_dir.mkdir(parents=True, exist_ok=True)
    pages_dir = output_dir / "pages"
    media_dir = output_dir / "media"
    pages_dir.mkdir(exist_ok=True)
    media_dir.mkdir(exist_ok=True)

    report = {
        "input": str(input_path),
        "output": str(output_dir),
        "slide_width": slide_w_px,
        "slide_height": slide_h_px,
        "slides": [],
        "summary": {
            "total_elements": 0,
            "high_confidence": 0,
            "medium_confidence": 0,
            "low_confidence": 0,
            "average_confidence": 0.0,
        },
    }

    manifest = {
        "version": "v2",
        "title": input_path.stem,
        "size": [ENGINE_W, ENGINE_H],
        "pages": [],
    }

    all_confidences = []

    for slide_idx, slide in enumerate(prs.slides):
        page_name = f"page_{slide_idx + 1:02d}.page"
        manifest["pages"].append(f"pages/{page_name}")

        page = {
            "pageType": "content",
            "elements": [],
        }

        slide_report = {
            "slide": slide_idx + 1,
            "elements": [],
        }

        for shape in slide.shapes:
            element = convert_shape(shape, slide_w, slide_h, media_dir, slide_idx)
            if element is None:
                continue

            confidence, reason = shape_confidence(
                element.get("_shape_type", "text"),
                element.get("_has_text", False),
                element.get("_has_fill", False),
                element.get("_is_group", False),
            )

            element["_confidence"] = confidence
            element["_reason"] = reason

            # 清理内部字段
            for key in list(element.keys()):
                if key.startswith("_"):
                    del element[key]

            page["elements"].append(element)

            slide_report["elements"].append({
                "elementId": element.get("elementId", ""),
                "type": element.get("elementType", "unknown"),
                "confidence": confidence,
                "reason": reason,
            })

            all_confidences.append(confidence)
            if confidence >= 0.9:
                report["summary"]["high_confidence"] += 1
            elif confidence >= 0.7:
                report["summary"]["medium_confidence"] += 1
            else:
                report["summary"]["low_confidence"] += 1

        report["slides"].append(slide_report)

        # 写入 page 文件
        page_path = pages_dir / page_name
        page_path.write_text(dump_yaml(page), encoding="utf-8")

    # 写入 manifest
    manifest_path = output_dir / f"{input_path.stem}.pptd"
    manifest_path.write_text(dump_yaml(manifest), encoding="utf-8")

    report["summary"]["total_elements"] = len(all_confidences)
    if all_confidences:
        report["summary"]["average_confidence"] = round(sum(all_confidences) / len(all_confidences), 4)

    return report


def convert_shape(shape, slide_w_px: float, slide_h_px: float, media_dir: Path, slide_idx: int) -> Optional[Dict[str, Any]]:
    """转换单个 shape 为 PPTD element。"""
    try:
        # 组合形状
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return convert_group(shape, slide_w_px, slide_h_px, media_dir, slide_idx)

        # 图片
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return convert_picture(shape, slide_w_px, slide_h_px, media_dir, slide_idx)

        # 表格
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return convert_table(shape, slide_w_px, slide_h_px)

        # 图表
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return convert_chart(shape, slide_w_px, slide_h_px)

        # 文本框 / 形状（有实际文本才按文本处理）
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                element = convert_text_shape(shape, slide_w_px, slide_h_px)
                if element:
                    return element

        # 纯形状（AUTO_SHAPE，无论有无空文本框）
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return convert_auto_shape(shape, slide_w_px, slide_h_px)

        # 连接线
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return convert_line(shape, slide_w_px, slide_h_px)

        # 其他：占位
        return {
            "elementId": f"slide{slide_idx + 1}_el_{uuid.uuid4().hex[:8]}",
            "elementType": "text",
            "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
            "content": {
                "text": f"[未转换: {shape.shape_type}]",
                "fontSize": 14,
                "color": "#FF0000",
            },
            "_shape_type": str(shape.shape_type),
            "_has_text": False,
            "_has_fill": False,
            "_is_group": False,
        }
    except Exception as e:
        return {
            "elementId": f"slide{slide_idx + 1}_el_{uuid.uuid4().hex[:8]}",
            "elementType": "text",
            "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
            "content": {
                "text": f"[转换失败: {e}]",
                "fontSize": 12,
                "color": "#FF0000",
            },
            "_shape_type": "error",
            "_has_text": False,
            "_has_fill": False,
            "_is_group": False,
        }


def convert_bounds(shape, slide_w_px: float, slide_h_px: float) -> List[float]:
    """shape 边界 → 引擎坐标。shape.left 等是 EMU（Length 对象），
    slide_w_px/h_px 也是 EMU（prs.slide_width），统一用 EMU 除后转引擎坐标。"""
    # shape.left 等返回 Length（EMU 整数），slide_w_px 传入的已经是 px
    # 但调用处传的 slide_w_px 实际上是 prs.slide_width（EMU），统一按 EMU 处理
    x = int(shape.left) / slide_w_px * ENGINE_W if shape.left else 0
    y = int(shape.top) / slide_h_px * ENGINE_H if shape.top else 0
    w = int(shape.width) / slide_w_px * ENGINE_W if shape.width else 0
    h = int(shape.height) / slide_h_px * ENGINE_H if shape.height else 0
    return [round(x, 2), round(y, 2), round(w, 2), round(h, 2)]


def convert_text_shape(shape, slide_w_px: float, slide_h_px: float) -> Optional[Dict[str, Any]]:
    """转换文本框/形状。"""
    if not shape.has_text_frame:
        return None

    text = shape.text_frame.text.strip()
    bounds = convert_bounds(shape, slide_w_px, slide_h_px)

    element = {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "text",
        "bounds": bounds,
        "content": {
            "text": text,
            "fontSize": 18,
            "color": "#000000",
        },
        "_shape_type": "text",
        "_has_text": bool(text),
        "_has_fill": shape.fill.type is not None,
        "_is_group": False,
    }

    # 提取字体信息
    if shape.text_frame.paragraphs:
        para = shape.text_frame.paragraphs[0]
        if para.runs:
            run = para.runs[0]
            try:
                if run.font.size:
                    # font.size 是 Length 对象，.pt 属性给磅值
                    element["content"]["fontSize"] = round(run.font.size.pt, 1)
            except (AttributeError, TypeError):
                pass
            try:
                if run.font.bold:
                    element["content"]["bold"] = True
                if run.font.italic:
                    element["content"]["italic"] = True
            except (AttributeError, TypeError):
                pass
            try:
                if run.font.color and run.font.color.rgb:
                    element["content"]["color"] = normalize_color(run.font.color.rgb)
            except (AttributeError, TypeError):
                pass

    # 对齐
    if para.alignment is not None:
        from pptx.enum.text import PP_ALIGN
        align_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
        }
        if para.alignment in align_map:
            element["content"]["align"] = [align_map[para.alignment], "top"]

    return element


def convert_auto_shape(shape, slide_w_px: float, slide_h_px: float) -> Dict[str, Any]:
    """转换预设形状。"""
    shape_name = "rect"
    try:
        # python-pptx: shape.auto_shape_type 返回 MSO_SHAPE 枚举
        from pptx.enum.shapes import MSO_SHAPE
        auto_type = shape.auto_shape_type
        if auto_type is not None:
            # MSO_SHAPE 枚举名 → PPTD shapeName
            # 枚举名如 RECTANGLE, ROUNDED_RECTANGLE, OVAL 等
            enum_name = str(auto_type).split(".")[-1].split(" ")[0].upper() if auto_type else ""
            ooxml_map = {
                "RECTANGLE": "rect", "ROUNDED_RECTANGLE": "roundRect",
                "OVAL": "ellipse", "ISOSCELES_TRIANGLE": "triangle",
                "RIGHT_TRIANGLE": "triangle", "DIAMOND": "diamond",
                "PARALLELOGRAM": "parallelogram", "TRAPEZOID": "trapezoid",
                "HEXAGON": "hexagon", "PENTAGON": "pentagon", "OCTAGON": "octagon",
                "RIGHT_ARROW": "rightArrow", "LEFT_ARROW": "leftArrow",
                "UP_ARROW": "upArrow", "DOWN_ARROW": "downArrow",
                "CHEVRON": "chevron", "HOME_PLATE": "homePlate",
                "STAR_5_POINT": "star5", "STAR_4_POINT": "star4", "STAR_6_POINT": "star6",
                "CROSS": "cross", "DONUT": "donut", "CLOUD": "cloud",
                "HEART": "heart", "LIGHTNING_BOLT": "lightningBolt",
                "SMILEY_FACE": "smileyFace", "SUN": "sun", "MOON": "moon",
                "BANNER": "banner", "RIBBON": "ribbon", "SEAL": "seal",
                "GEAR_6": "gear", "FOLDED_CORNER": "foldedCorner",
                "PIE": "pie", "ARC": "arc", "BLOCK_ARC": "blockArc",
            }
            if enum_name in ooxml_map:
                shape_name = ooxml_map[enum_name]
            else:
                # 尝试从 MSO_SHAPE 枚举成员名反向查找
                for name, val in MSO_SHAPE.__members__.items():
                    if val == auto_type:
                        shape_name = OOXML_TO_PPTD_SHAPE.get(name, name.lower())
                        break
    except Exception:
        pass

    element = {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "shape",
        "shapeName": shape_name,
        "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
        "_shape_type": shape_name,
        "_has_text": False,
        "_has_fill": shape.fill.type is not None,
        "_is_group": False,
    }

    # 填充
    try:
        if shape.fill.type is not None:
            if shape.fill.fore_color and shape.fill.fore_color.rgb:
                element["fill"] = {"type": "solid", "color": normalize_color(shape.fill.fore_color.rgb)}
    except (AttributeError, TypeError):
        pass

    return element


def convert_picture(shape, slide_w_px: float, slide_h_px: float, media_dir: Path, slide_idx: int) -> Dict[str, Any]:
    """转换图片。"""
    image = shape.image
    ext = image.ext or "png"
    filename = f"slide{slide_idx + 1}_img_{uuid.uuid4().hex[:8]}.{ext}"
    media_path = media_dir / filename
    media_path.write_bytes(image.blob)

    return {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "image",
        "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
        "src": f"media/{filename}",
        "fit": {"mode": "cover"},
        "_shape_type": "picture",
        "_has_text": False,
        "_has_fill": False,
        "_is_group": False,
    }


def convert_table(shape, slide_w_px: float, slide_h_px: float) -> Dict[str, Any]:
    """转换表格。"""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append({
                "text": cell.text.strip(),
                "colSpan": cell._tc.gridSpan if hasattr(cell._tc, 'gridSpan') else 1,
                "rowSpan": cell._tc.rowSpan if hasattr(cell._tc, 'rowSpan') else 1,
            })
        rows.append(cells)

    return {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "table",
        "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
        "rows": rows,
        "columnWidths": [1.0 / len(table.columns)] * len(table.columns),
        "_shape_type": "table",
        "_has_text": True,
        "_has_fill": False,
        "_is_group": False,
    }


def convert_chart(shape, slide_w_px: float, slide_h_px: float) -> Dict[str, Any]:
    """转换图表（简化：提取数据）。"""
    chart = shape.chart
    chart_type = "bar"
    try:
        from pptx.enum.chart import XL_CHART_TYPE
        type_map = {
            XL_CHART_TYPE.COLUMN_CLUSTERED: "bar",
            XL_CHART_TYPE.BAR_CLUSTERED: "bar",
            XL_CHART_TYPE.LINE: "line",
            XL_CHART_TYPE.PIE: "pie",
            XL_CHART_TYPE.DOUGHNUT: "donut",
            XL_CHART_TYPE.RADAR: "radar",
            XL_CHART_TYPE.XY_SCATTER: "scatter",
        }
        for enum_val, name in type_map.items():
            if chart.chart_type == enum_val:
                chart_type = name
                break
    except Exception:
        pass

    # 提取数据
    categories = []
    rows = []
    try:
        for i, cat in enumerate(chart.plots[0].categories):
            categories.append(str(cat))
        for si, series in enumerate(chart.series):
            values = [v for v in series.values]
            rows.append(values)
    except Exception:
        pass

    cols = ["Category"] + [f"Series {i+1}" for i in range(len(rows))]

    return {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "chart",
        "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
        "data": {
            "cols": cols,
            "rows": [([categories[i]] + [r[i] for r in rows]) for i in range(len(categories))],
        },
        "series": [{"type": chart_type, "encode": {"x": cols[0], "y": cols[1]}}],
        "_shape_type": "chart",
        "_has_text": False,
        "_has_fill": False,
        "_is_group": False,
    }


def convert_line(shape, slide_w_px: float, slide_h_px: float) -> Dict[str, Any]:
    """转换连接线。"""
    bounds = convert_bounds(shape, slide_w_px, slide_h_px)
    x, y, w, h = bounds
    return {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "line",
        "bounds": bounds,
        "points": [
            [0, 0],
            [w, h],
        ],
        "border": {"color": "#000000", "width": 1},
        "_shape_type": "line",
        "_has_text": False,
        "_has_fill": False,
        "_is_group": False,
    }


def convert_group(shape, slide_w_px: float, slide_h_px: float, media_dir: Path, slide_idx: int) -> Dict[str, Any]:
    """转换组合形状（简化：转为占位）。"""
    return {
        "elementId": f"el_{uuid.uuid4().hex[:8]}",
        "elementType": "text",
        "bounds": convert_bounds(shape, slide_w_px, slide_h_px),
        "content": {
            "text": "[组合形状：需解组后逐元素处理]",
            "fontSize": 12,
            "color": "#FF0000",
        },
        "_shape_type": "group",
        "_has_text": False,
        "_has_fill": False,
        "_is_group": True,
    }


def dump_yaml(obj: Any) -> str:
    """简单 YAML 序列化（不依赖 pyyaml）。"""
    return _yaml_dump(obj, 0) + "\n"


def _yaml_dump(obj: Any, indent: int) -> str:
    prefix = "  " * indent
    if obj is None:
        return "null"
    if isinstance(obj, bool):
        return "true" if obj else "false"
    if isinstance(obj, int):
        return str(obj)
    if isinstance(obj, float):
        return f"{obj:.2f}" if obj != int(obj) else str(int(obj))
    if isinstance(obj, str):
        if any(c in obj for c in [":", "#", "\n", "'", '"', "{", "[", ",", "&", "*", "?", "|", "<", ">", "=", "!", "%", "@", "`"]):
            escaped = obj.replace("\\", "\\\\").replace('"', '\\"').replace("\n", "\\n")
            return f'"{escaped}"'
        return obj
    if isinstance(obj, list):
        if not obj:
            return "[]"
        lines = []
        for item in obj:
            item_str = _yaml_dump(item, indent + 1)
            if isinstance(item, dict) and item:
                # dict in list: "- key: value" 同行，子项缩进
                first = True
                for line in item_str.split("\n"):
                    if first:
                        lines.append(f"{prefix}- {line}")
                        first = False
                    else:
                        lines.append(f"  {prefix}{line}")
            elif isinstance(item, list) and item:
                # list in list: "- - ..." 同行
                first = True
                for line in item_str.split("\n"):
                    if first:
                        lines.append(f"{prefix}- {line}")
                        first = False
                    else:
                        lines.append(f"  {prefix}{line}")
            else:
                lines.append(f"{prefix}- {item_str}")
        return "\n".join(lines)
    if isinstance(obj, dict):
        if not obj:
            return "{}"
        lines = []
        for k, v in obj.items():
            v_str = _yaml_dump(v, indent + 1)
            if isinstance(v, (dict, list)) and v:
                lines.append(f"{prefix}{k}:\n{v_str}")
            else:
                lines.append(f"{prefix}{k}: {v_str}")
        return "\n".join(lines)
    return str(obj)


def main():
    ap = argparse.ArgumentParser(description="PPTX → PPTD 转换 + 保真度评分")
    ap.add_argument("input", help="输入 PPTX 文件路径")
    ap.add_argument("output", help="输出 PPTD 项目目录")
    ap.add_argument("--json", action="store_true", help="仅输出 JSON 报告")
    ap.add_argument("--report", type=Path, help="报告输出路径（默认: <output>/fidelity-report.json）")
    args = ap.parse_args()

    input_path = Path(args.input).expanduser().resolve()
    output_dir = Path(args.output).expanduser().resolve()

    if not input_path.exists():
        print(f"error: 输入文件不存在: {input_path}", file=sys.stderr)
        return 1

    if not input_path.suffix.lower() == ".pptx":
        print(f"error: 输入必须是 .pptx 文件: {input_path}", file=sys.stderr)
        return 1

    report = convert_pptx_to_pptd(input_path, output_dir)

    report_path = args.report or (output_dir / "fidelity-report.json")
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        s = report["summary"]
        print(f"convert_fidelity: {input_path.name} → {output_dir}")
        print(f"  元素总数: {s['total_elements']}")
        print(f"  高置信度 (≥90%): {s['high_confidence']}")
        print(f"  中置信度 (70-90%): {s['medium_confidence']}")
        print(f"  低置信度 (<70%): {s['low_confidence']}")
        print(f"  平均置信度: {s['average_confidence']:.1%}")
        print(f"  报告: {report_path}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
